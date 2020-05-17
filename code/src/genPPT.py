from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer

#total number of slides excluding the first slide
global num_slides


def init(title_text, subtitle_text):
    ppt = Presentation()
    title_slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    if len(subtitle_text) != 0:
        subtitle.text = subtitle_text

    add_logo(slide)
    return ppt

def add_logo(slide):
    img_path = 'static/images/logo.jpg'
    left = Inches(2)
    pic = slide.shapes.add_picture(img_path, left, 0)

def summarize(file_path, num_slides = 10, points_per_slide = 3):
    parser = PlaintextParser.from_file(file_path, Tokenizer("english"))
    summarizer = LexRankSummarizer()
    summary = summarizer(parser.document, points_per_slide * num_slides)

    top_lines = []

    for sentence in summary:
        top_lines.append(str(sentence))
    return top_lines

def create_slide(ppt, sentences):
    bullet_slide_layout = ppt.slide_layouts[1]
    slide = ppt.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Topic Name'

    tf = body_shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    
    for sentence in sentences:
        p = tf.add_paragraph()
        p.font.size = Pt(18)     
        p.text = sentence
        p.level = 0
   
#returns a list of sentences list of respective slide
#format [[slide1SentencesList, ...]]
def get_ppt_points(top_sentences, points_per_slide):
    global num_slides
    ppt_points = []
    for slide_no in range(num_slides):
        # not enough sentences to write to a new slide
        if (len(top_sentences) - (points_per_slide * (slide_no + 1))) < points_per_slide:
            break
        ppt_points.append([])

        #add points of a particular slide to a sperate lists
        for point_num in range(points_per_slide):
            ppt_points[-1].append(top_sentences[points_per_slide * slide_no + point_num])
    return ppt_points


def createPPT(title, file_path, subtitle = '', no_slides = 5, points_per_slide = 3):
    global num_slides
    num_slides = no_slides
    ppt = init(title, subtitle)
    sentences = summarize(file_path, num_slides, points_per_slide)
    ppt_points = get_ppt_points(sentences, points_per_slide)

    for slide_points in ppt_points:
        if len(slide_points) == 0:
            continue
        create_slide(ppt, slide_points)

    ppt.save('notes.pptx')


#createPPT('notes.txt', 'notes.txt', 'title')